"""
Generate architecture comparison PPTX for MioServer / CodeIsland / CodeLight.
Covers 4 deployment scenarios with data flow diagrams.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────────────────────────────
BRAND_LIME  = RGBColor(0xCA, 0xFF, 0x00)   # #CAFF00 MioIsland brand
BRAND_DARK  = RGBColor(0x1A, 0x1A, 0x1A)   # near-black
BRAND_GRAY  = RGBColor(0x2D, 0x2D, 0x2D)   # card bg
BRAND_MID   = RGBColor(0x50, 0x50, 0x50)   # subtitle text
TEXT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIME   = BRAND_LIME
TEXT_DIM    = RGBColor(0x99, 0x99, 0x99)
ACCENT_BLUE = RGBColor(0x4A, 0x9D, 0xFF)
ACCENT_RED  = RGBColor(0xFF, 0x4A, 0x4A)
ACCENT_GRN  = RGBColor(0x4A, 0xFF, 0x8A)
LINE_COLOR  = RGBColor(0x60, 0x60, 0x60)
BG_DARK     = RGBColor(0x12, 0x12, 0x12)

# ── Helpers ──────────────────────────────────────────────────────────────────

def blank_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    return sl


def rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    """Add a filled rectangle. All units in inches."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w or 1)
    else:
        shape.line.fill.background()
    return shape


def rounded_rect(slide, l, t, w, h, fill, corner=0.1, line=None):
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def label(slide, text, l, t, w, h=None, size=12, bold=False,
          color=TEXT_WHITE, align=PP_ALIGN.CENTER, wrap=True):
    """Add a text box. h defaults to 0.4 if not given."""
    if h is None:
        h = max(0.3, round(size * 0.06, 2))
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def subtitle(slide, text, top=0.3, color=TEXT_DIM):
    return label(slide, text, 0.3, top, 9.4, size=13, color=color)


def page_num(slide, n, total):
    label(slide, f"{n} / {total}", 0.2, 8.8, 1.5, size=9, color=TEXT_DIM,
          align=PP_ALIGN.LEFT)


def arrow(slide, x1, y1, x2, y2, color=LINE_COLOR, w=1.5, dash=False):
    """Draw a simple line connector. For arrows we add a triangle marker."""
    from pptx.util import Pt
    from pptx.enum.dml import MSO_THEME_COLOR
    conn = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    return conn


def add_text_box(slide, text, l, t, w, h, font_size=11, bold=False,
                  color=TEXT_WHITE, align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


# ── Node box helper ─────────────────────────────────────────────────────────
# Each node: (label, sub_label, fill, text_color)
NODE_H = 0.7
NODE_W = 1.8

def node(slide, text, sub, cx, cy, fill=BRAND_GRAY, tcolor=TEXT_WHITE,
         w=NODE_W, h=NODE_H):
    """Draw a rounded-rect node centered at (cx, cy)."""
    l = cx - w / 2
    t = cy - h / 2
    box = rounded_rect(slide, l, t, w, h, fill, corner=0.08, line=LINE_COLOR)
    label(slide, text, l + 0.05, t + 0.05, w - 0.1, 0.32, size=11, bold=True,
          color=tcolor)
    if sub:
        label(slide, sub, l + 0.05, t + 0.32, w - 0.1, 0.3, size=8,
              color=TEXT_DIM)
    return (l, t, w, h)


def dashed_line(slide, x1, y1, x2, y2, color=BRAND_LIME, w=1.2):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    import copy
    from lxml import etree
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    # Make dashed via XML
    ln = conn.line._ln
    prst = ln.find(qn('a:prstDash'))
    if prst is None:
        prst = etree.SubElement(ln, qn('a:prstDash'))
    prst.set('val', 'dash')
    return conn


def solid_line(slide, x1, y1, x2, y2, color=LINE_COLOR, w=1.2):
    from pptx.util import Pt
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    return conn


def arrowhead_line(slide, x1, y1, x2, y2, color=LINE_COLOR, w=1.2):
    """Line with arrow marker."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    # Add arrow head
    tailEnd = conn.line._ln
    # Create <a:tailEnd> element
    tail_end = etree.SubElement(tailEnd, qn('a:tailEnd'))
    tail_end.set('type', 'triangle')
    tail_end.set('w', 'med')
    tail_end.set('len', 'med')
    return conn


# ── Box with icon + text ──────────────────────────────────────────────────────
def component_box(slide, icon_text, title, subtitle_text, cx, cy,
                  fill=BRAND_GRAY, title_color=TEXT_WHITE,
                  sub_color=TEXT_DIM, w=1.7, h=0.85):
    l = cx - w / 2
    t = cy - h / 2
    rounded_rect(slide, l, t, w, h, fill, corner=0.08, line=LINE_COLOR)
    # Icon (text emoji substitute)
    label(slide, icon_text, l, t + 0.06, w, 0.32, size=16, bold=True,
          color=BRAND_LIME)
    label(slide, title, l, t + 0.38, w, 0.22, size=9, bold=True,
          color=title_color)
    label(slide, subtitle_text, l, t + 0.56, w, 0.22, size=7.5,
          color=sub_color)


# ── Flow label ────────────────────────────────────────────────────────────────
def flow_tag(slide, text, cx, cy, fill=BRAND_GRAY, color=TEXT_LIME, w=1.4, h=0.28):
    l = cx - w / 2
    t = cy - h / 2
    rounded_rect(slide, l, t, w, h, fill, corner=0.05, line=color)
    label(slide, text, l, t, w, h, size=7.5, bold=True, color=color)


# ── Legend ───────────────────────────────────────────────────────────────────
def legend(slide, items, l, t):
    """items: list of (color, label)"""
    for i, (col, lbl) in enumerate(items):
        y = t + i * 0.28
        rect(slide, l, y, 0.18, 0.14, fill=col)
        label(slide, lbl, l + 0.22, y - 0.04, 1.5, 0.22, size=8,
              color=TEXT_DIM)


# ── Slide builder functions ──────────────────────────────────────────────────

def slide_title(prs, title, subtitle=None):
    sl = blank_slide(prs)
    # Background
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    # Accent bar top
    rect(sl, 0, 0, 10, 0.08, fill=BRAND_LIME)
    rect(sl, 0, 8.42, 10, 0.08, fill=BRAND_LIME)
    # Title
    label(sl, title, 0.5, 3.0, 9, size=36, bold=True, color=TEXT_WHITE)
    if subtitle:
        label(sl, subtitle, 0.5, 3.85, 9, size=16, color=TEXT_DIM)
    return sl


def slide_section(prs, title, section_num=None):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    rect(sl, 0, 0, 10, 0.06, fill=BRAND_LIME)
    if section_num:
        label(sl, section_num, 0.4, 3.4, 1, size=72, bold=True,
              color=BRAND_LIME, align=PP_ALIGN.LEFT)
    label(sl, title, 0.4 + (1.2 if section_num else 0), 3.5,
          8.5, size=32, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT)
    return sl


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════════════════
def cover(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    rect(sl, 0, 0, 10, 0.1, fill=BRAND_LIME)
    rect(sl, 0, 8.4, 10, 0.1, fill=BRAND_LIME)

    # Left accent column
    rect(sl, 0, 0.1, 0.12, 8.3, fill=BRAND_LIME)

    # Title
    label(sl, "MioIsland × MioServer", 0.5, 1.8, 9, size=40, bold=True,
          color=TEXT_WHITE)
    label(sl, "框架与链路对比", 0.5, 2.6, 9, size=28, bold=False,
          color=BRAND_LIME)
    label(sl, "Architecture & Data Flow Comparison", 0.5, 3.1, 9, size=14,
          color=TEXT_DIM)

    # Divider
    rect(sl, 0.5, 3.6, 5, 0.03, fill=LINE_COLOR)

    # Topics list
    topics = [
        "CodeIsland / MioIsland / CodeLight 三者关系",
        "4 种部署架构详解",
        "数据流向与 session 生命周期",
        "Hook → mio-server 接入方案",
        "iPhone 配对与消息推送链路",
    ]
    for i, t in enumerate(topics):
        label(sl, f"›  {t}", 0.5, 3.8 + i * 0.44, 8.5, size=13,
              color=TEXT_WHITE)

    label(sl, "2026-04-30", 0.5, 8.1, 3, size=10, color=TEXT_DIM)
    return sl


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Three Musketeers
# ═══════════════════════════════════════════════════════════════════════════════
def slide_overview(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    rect(sl, 0, 0, 10, 0.06, fill=BRAND_LIME)

    label(sl, "01  三组件定位", 0.4, 0.25, 9, size=22, bold=True,
          color=TEXT_WHITE)
    subtitle(sl, "CodeIsland · MioIsland · CodeLight — 各司其职", top=0.75)

    cards = [
        ("🖥", "CodeIsland",      "Mac notch 桌面应用",     BRAND_GRAY,
         "• SSH relay 远程 CC\n• 本地 session 读取\n• tmux 进程管理\n• Notch UI 显示"),
        ("📱", "MioIsland",       "iPhone 控制端（经典）",  RGBColor(0x20, 0x30, 0x20),
         "• SSH 连接 Mac\n• 远程控制 CC\n• 接收推送通知\n• Launch Preset 管理"),
        ("📲", "CodeLight",       "iPhone 控制端（新）",    RGBColor(0x15, 0x15, 0x30),
         "• 直连 mio-server\n• 扫码 / shortCode 配对\n• 接收 APNs 通知\n• 不需要 Mac 在线"),
    ]

    for i, (icon, title, sub, fill, bullets) in enumerate(cards):
        x = 0.4 + i * 3.15
        # Card bg
        rounded_rect(sl, x, 1.35, 3.0, 5.2, fill, corner=0.1,
                     line=LINE_COLOR)
        # Header strip
        rect(sl, x, 1.35, 3.0, 0.7, fill=BRAND_LIME)
        # Icon
        label(sl, icon, x, 1.38, 3.0, 0.6, size=24, bold=True,
              color=BRAND_DARK)
        # Title
        label(sl, title, x, 2.1, 3.0, 0.35, size=15, bold=True,
              color=TEXT_WHITE)
        label(sl, sub, x, 2.42, 3.0, 0.28, size=9, color=TEXT_DIM)

        # Bullets
        txb = sl.shapes.add_textbox(
            Inches(x + 0.18), Inches(2.78), Inches(2.7), Inches(3.5))
        tf = txb.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets.split('\n')):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = b
            run.font.size = Pt(10)
            run.font.color.rgb = TEXT_DIM

    # Bottom note
    rect(sl, 0.4, 6.75, 9.2, 0.03, fill=LINE_COLOR)
    label(sl, "关键区别：CodeLight 不走 SSH 直连 Mac，而是通过 mio-server 作为消息中继", 0.4, 6.85, 9.2, size=10,
          color=BRAND_LIME)
    return sl


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Current Architecture
# ═══════════════════════════════════════════════════════════════════════════════
def slide_arch_current(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    rect(sl, 0, 0, 10, 0.06, fill=BRAND_LIME)

    label(sl, "02  现有架构", 0.4, 0.25, 9, size=22, bold=True,
          color=TEXT_WHITE)
    subtitle(sl, "CC 运行在本地 Mac，通过 hook socket 与 notch 联动", top=0.75)

    # Legend
    legend(sl, [
        (BRAND_LIME,  "实时数据流（Socket.io）"),
        (ACCENT_BLUE, "消息轮询（HTTP）"),
        (ACCENT_RED,  "SSH 隧道"),
        (LINE_COLOR,  "间接关联"),
    ], 0.4, 1.2)

    # ── Nodes ────────────────────────────────────────────────────────────────
    CY = 3.6
    GAP = 2.2

    # iPhone
    node(sl, "iPhone", "CodeLight App",
         0.8, CY, fill=RGBColor(0x15, 0x15, 0x30), tcolor=TEXT_WHITE, w=1.5)

    # mio-server
    node(sl, "mio-server", "Railway · Socket.io",
         0.8 + GAP, CY, fill=RGBColor(0x0A, 0x20, 0x0A), tcolor=BRAND_LIME, w=1.8)

    # Mac notch
    node(sl, "MioIsland", "Mac notch app",
         0.8 + GAP * 2, CY, fill=BRAND_GRAY, tcolor=TEXT_WHITE, w=1.6)

    # Claude Code
    node(sl, "Claude Code", "本地运行",
         0.8 + GAP * 3, CY, fill=RGBColor(0x20, 0x0A, 0x0A),
         tcolor=ACCENT_RED, w=1.6)

    # Hook Socket
    node(sl, "Hook Socket", "127.0.0.1:9871",
         0.8 + GAP * 2, CY + 1.4, fill=RGBColor(0x30, 0x20, 0x10),
         tcolor=ACCENT_BLUE, w=1.9, h=0.6)

    # ── Arrows ──────────────────────────────────────────────────────────────
    # iPhone ↔ mio-server
    dashed_line(sl, 2.0, CY, 0.8 + GAP - 1.0, CY, color=BRAND_LIME, w=1.5)
    flow_tag(sl, "Socket.io /v1/updates", 1.4, CY - 0.22,
             fill=RGBColor(0x10, 0x20, 0x10), color=BRAND_LIME, w=1.9)

    # mio-server ↔ Mac notch
    solid_line(sl, 0.8 + GAP + 1.0, CY, 0.8 + GAP * 2 - 1.0, CY,
               color=BRAND_LIME, w=1.5)
    flow_tag(sl, "Socket.io /v1/updates", 0.8 + GAP + 0.5, CY - 0.22,
             fill=RGBColor(0x10, 0x20, 0x10), color=BRAND_LIME, w=1.9)

    # Mac notch ↔ hook socket
    solid_line(sl, 0.8 + GAP * 2, CY + 0.38, 0.8 + GAP * 2, CY + 1.0,
               color=ACCENT_BLUE, w=1.2)
    flow_tag(sl, "hook socket", 0.8 + GAP * 2 + 0.12, CY + 0.7,
             fill=RGBColor(0x10, 0x10, 0x20), color=ACCENT_BLUE, w=1.1)

    # Hook ↔ Claude Code
    solid_line(sl, 0.8 + GAP * 2 + 1.0, CY + 1.4,
               0.8 + GAP * 3 - 1.0, CY,
               color=LINE_COLOR, w=1.0)

    # ── Summary ──────────────────────────────────────────────────────────────
    rect(sl, 0.4, 6.0, 9.2, 1.85, fill=RGBColor(0x0D, 0x0D, 0x0D),
         line=LINE_COLOR)
    label(sl, "数据流", 0.6, 6.08, 1.2, size=10, bold=True, color=BRAND_LIME)
    rows = [
        ("1", "iPhone 扫码配对 → mio-server 建立 DeviceLink"),
        ("2", "CC hook 写入 JSONL → MioIsland 读取并通过 Socket.io 发送到 mio-server"),
        ("3", "mio-server 存储 session 消息 → iPhone 轮询 /v1/sessions 获取更新"),
        ("4", "iPhone 也可以发消息 → mio-server → MioIsland → 本地 CC 终端"),
    ]
    for i, (num, desc) in enumerate(rows):
        label(sl, num, 0.6 + i * 0, 6.32 + i * 0.34, 0.3, size=9,
              bold=True, color=BRAND_LIME)
        label(sl, desc, 0.9, 6.32 + i * 0.34, 8.5, size=9, color=TEXT_DIM)
    return sl


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Server CC Architecture
# ═══════════════════════════════════════════════════════════════════════════════
def slide_arch_server(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    rect(sl, 0, 0, 10, 0.06, fill=BRAND_LIME)

    label(sl, "03  架构 A", 0.4, 0.25, 9, size=22, bold=True,
          color=TEXT_WHITE)
    subtitle(sl, "CC 运行在远程服务器，hook 直连 mio-server · 不需要 Mac 在线", top=0.75)

    CY = 2.9
    WIDTHS = [1.5, 2.0, 1.5, 1.8]
    CX = [0.8, 2.9, 5.5, 7.8]
    FILLS = [
        RGBColor(0x15, 0x15, 0x30),   # iPhone
        RGBColor(0x0A, 0x20, 0x0A),   # mio-server
        RGBColor(0x20, 0x0A, 0x0A),   # CC server
        RGBColor(0x10, 0x10, 0x20),   # hook client
    ]
    TCOLORS = [TEXT_WHITE, BRAND_LIME, ACCENT_RED, ACCENT_BLUE]

    labels_data = [
        ("iPhone", "CodeLight App"),
        ("mio-server", "Railway · 已部署"),
        ("Claude Code", "远程服务器"),
        ("Hook Client", "hook socket → mio-server"),
    ]
    for i, (lbl, sub) in enumerate(labels_data):
        node(sl, lbl, sub, CX[i], CY, fill=FILLS[i], tcolor=TCOLORS[i],
             w=WIDTHS[i])

    # Hook client below CC
    node(sl, "Hook Socket", "127.0.0.1:port",
         CX[2], CY + 1.35, fill=RGBColor(0x30, 0x20, 0x10),
         tcolor=ACCENT_BLUE, w=1.9, h=0.55)

    # ── Arrows ──────────────────────────────────────────────────────────────
    # iPhone ↔ mio-server
    dashed_line(sl, CX[0] + 0.78, CY, CX[1] - 1.1, CY, color=BRAND_LIME, w=1.5)
    flow_tag(sl, "Socket.io /v1/updates", (CX[0] + CX[1]) / 2, CY - 0.22,
             fill=RGBColor(0x10, 0x20, 0x10), color=BRAND_LIME, w=2.0)

    # mio-server ↔ CC server
    dashed_line(sl, CX[1] + 1.1, CY, CX[2] - 0.85, CY, color=BRAND_LIME, w=1.5)
    flow_tag(sl, "Socket.io /v1/updates", (CX[1] + CX[2]) / 2, CY - 0.22,
             fill=RGBColor(0x10, 0x20, 0x10), color=BRAND_LIME, w=2.0)

    # CC ↔ hook
    solid_line(sl, CX[2], CY + 0.38, CX[2], CY + 1.0,
               color=LINE_COLOR, w=1.0)

    # ── Problem highlight ─────────────────────────────────────────────────────
    rect(sl, 0.4, 5.1, 9.2, 2.85, fill=RGBColor(0x0D, 0x0D, 0x0D),
         line=LINE_COLOR)

    label(sl, "⚠  核心挑战", 0.6, 5.2, 3, size=11, bold=True, color=ACCENT_RED)

    challenges = [
        ("DeviceId 问题",
         "session 归属哪个 deviceId？mio-server 的 session API 绑定 Mac deviceId。"
         "Server CC 需要冒充 Mac deviceId 或新建独立 session 体系。"),
        ("HTTP API 替代 Socket.io",
         "hook 不支持 WebSocket，可以直接 POST /v1/sessions/:id/messages。"
         "需要 hook socket 改为 HTTP 客户端。"),
        ("无 Mac 时 deviceLink",
         "iPhone 配对的是 Mac shortCode。Server CC 模式需要建立新的配对机制。"),
    ]
    for i, (title, desc) in enumerate(challenges):
        label(sl, f"› {title}", 0.6, 5.48 + i * 0.78, 2.5, size=9,
              bold=True, color=BRAND_LIME)
        label(sl, desc, 0.6, 5.72 + i * 0.78, 8.8, size=8.5, color=TEXT_DIM)

    return sl


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SSH Relay Architecture
# ═══════════════════════════════════════════════════════════════════════════════
def slide_arch_ssh(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    rect(sl, 0, 0, 10, 0.06, fill=BRAND_LIME)

    label(sl, "04  架构 B", 0.4, 0.25, 9, size=22, bold=True,
          color=TEXT_WHITE)
    subtitle(sl, "SSH relay 模式 · Mac notch 作为 SSH client · 远程 CC session 映射到 notch", top=0.75)

    # Draw the topology (7 nodes in two rows)
    CY1 = 2.6
    CY2 = 5.0
    CX_LEFT  = 1.0
    CX_MID   = 3.8
    CX_RIGHT = 6.8

    node(sl, "iPhone", "CodeLight", CX_LEFT, CY1,
         fill=RGBColor(0x15, 0x15, 0x30), w=1.4)
    node(sl, "mio-server", "Railway", CX_MID, CY1,
         fill=RGBColor(0x0A, 0x20, 0x0A), tcolor=BRAND_LIME, w=1.8)
    node(sl, "Mac Notch", "MioIsland", CX_RIGHT, CY1,
         fill=BRAND_GRAY, w=1.5)

    node(sl, "SSH Relay", "127.0.0.1:9871", CX_MID, CY2,
         fill=RGBColor(0x30, 0x20, 0x10), tcolor=ACCENT_BLUE, w=1.8)
    node(sl, "CC (Remote)", "SSH 远程机器", CX_RIGHT, CY2,
         fill=RGBColor(0x20, 0x0A, 0x0A), tcolor=ACCENT_RED, w=1.7)

    # ── Arrows ──────────────────────────────────────────────────────────────
    # iPhone ↔ mio-server
    dashed_line(sl, CX_LEFT + 0.73, CY1, CX_MID - 1.0, CY1,
               color=BRAND_LIME, w=1.5)
    flow_tag(sl, "Socket.io", (CX_LEFT + CX_MID) / 2, CY1 - 0.22,
             fill=RGBColor(0x10, 0x20, 0x10), color=BRAND_LIME, w=1.2)

    # mio-server ↔ Mac Notch
    dashed_line(sl, CX_MID + 1.0, CY1, CX_RIGHT - 0.85, CY1,
               color=BRAND_LIME, w=1.5)
    flow_tag(sl, "Socket.io", (CX_MID + CX_RIGHT) / 2, CY1 - 0.22,
             fill=RGBColor(0x10, 0x20, 0x10), color=BRAND_LIME, w=1.2)

    # Mac Notch → SSH Relay
    solid_line(sl, CX_RIGHT, CY1 - 0.38, CX_RIGHT, CY2 + 0.38,
               color=ACCENT_RED, w=1.2)
    flow_tag(sl, "SSH tunnel", CX_RIGHT + 0.1, (CY1 + CY2) / 2,
             fill=RGBColor(0x20, 0x10, 0x10), color=ACCENT_RED, w=1.0)

    # SSH Relay ↔ Remote CC
    solid_line(sl, CX_MID + 1.0, CY2, CX_RIGHT - 0.88, CY2,
               color=LINE_COLOR, w=1.0)
    flow_tag(sl, "local 127.0.0.1", (CX_MID + CX_RIGHT) / 2, CY2 - 0.22,
             fill=RGBColor(0x15, 0x15, 0x15), color=TEXT_DIM, w=1.6)

    # Hook → SSH Relay
    solid_line(sl, CX_RIGHT, CY2 + 0.38, CX_RIGHT, CY2 + 0.9,
               color=ACCENT_BLUE, w=1.0)
    flow_tag(sl, "hook socket", CX_RIGHT + 0.1, CY2 + 0.65,
             fill=RGBColor(0x10, 0x10, 0x20), color=ACCENT_BLUE, w=1.1)

    # ── Key difference ───────────────────────────────────────────────────────
    rect(sl, 0.4, 6.6, 9.2, 1.4, fill=RGBColor(0x0D, 0x0D, 0x0D),
         line=LINE_COLOR)
    label(sl, "工作原理", 0.6, 6.68, 2, size=10, bold=True, color=BRAND_LIME)
    steps = [
        "Mac notch 主动 SSH 到远程机器，建立 port forward：本地 9871 ↔ 远程 9871",
        "远程 CC 的 hook 连接到本地 127.0.0.1:9871，数据经 SSH 隧道传回 notch",
        "notch 把 session 数据通过 Socket.io 转发给 mio-server，iPhone 正常轮询",
        "iPhone 发消息 → mio-server → notch → SSH tunnel → 远程 CC",
    ]
    for i, s in enumerate(steps):
        label(sl, f"› {s}", 0.6, 6.88 + i * 0.28, 8.8, size=8.5, color=TEXT_DIM)
    return sl


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Comparison Table
# ═══════════════════════════════════════════════════════════════════════════════
def slide_comparison(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    rect(sl, 0, 0, 10, 0.06, fill=BRAND_LIME)

    label(sl, "05  架构对比", 0.4, 0.25, 9, size=22, bold=True,
          color=TEXT_WHITE)
    subtitle(sl, "4 种部署模式一览", top=0.75)

    headers = ["维度", "现有架构\n(本地CC)", "架构 A\n(Server CC)", "架构 B\n(SSH Relay)", "架构 C\n(Hybrid)"]
    col_x   = [0.4, 2.2, 4.1, 6.1, 8.1]
    col_w   = 1.7
    ROWS = [
        ("CC 运行环境", "本地 Mac",         "远程服务器",       "远程机器 SSH",      "本地 + 远程混合"),
        ("Mac 在线？",  "✓ 必须",           "✗ 不需要",         "✓ 必须",            "✓ 必须（notch桥接）"),
        ("iPhone 连接", "mio-server",       "mio-server",       "mio-server",        "mio-server"),
        ("Hook 目标",   "本地 127.0.0.1",    "mio-server HTTP",  "本地 127.0.0.1\n(经SSH隧道)", "本地 + mio-server"),
        ("DeviceId",    "Mac deviceId",      "需新建体系",       "Mac deviceId",      "Mac deviceId"),
        ("Session 存储","mio-server DB",     "mio-server DB\n(需改造)", "mio-server DB",  "mio-server DB"),
        ("APNs 通知",   "✓ 支持",            "需新增机制",        "✓ 支持",            "✓ 支持"),
        ("部署难度",    "★☆☆☆☆",           "★★★☆☆",          "★★☆☆☆",           "★★★☆☆"),
        ("适合场景",    "日常开发",           "无 Mac 环境",       "远程开发",           "多机器协作"),
    ]

    HDR_H = 0.65
    ROW_H = 0.65
    TBL_TOP = 1.25

    # Header row
    for j, (hdr, x) in enumerate(zip(headers, col_x)):
        fill = BRAND_LIME if j == 0 else RGBColor(0x18, 0x18, 0x18)
        col = BRAND_DARK if j == 0 else TEXT_WHITE
        rounded_rect(sl, x, TBL_TOP, col_w, HDR_H, fill, corner=0.05,
                     line=LINE_COLOR)
        label(sl, hdr, x + 0.06, TBL_TOP + 0.06, col_w - 0.1, HDR_H - 0.1,
              size=9, bold=True, color=col)

    # Data rows
    for i, row in enumerate(ROWS):
        y = TBL_TOP + HDR_H + i * ROW_H
        fill = RGBColor(0x0F, 0x0F, 0x0F) if i % 2 == 0 else RGBColor(0x14, 0x14, 0x14)
        for j, (cell, x) in enumerate(zip(row, col_x)):
            cell_fill = BRAND_GRAY if j == 0 else fill
            cell_col  = BRAND_LIME if j == 0 else TEXT_WHITE
            # Highlight cells with ✓ or ✗
            if j > 0 and any(c in cell for c in ["✓", "★", "★★★"]):
                cell_col = ACCENT_GRN
            if j > 0 and any(c in cell for c in ["✗", "★★☆☆☆"]):
                cell_col = ACCENT_RED
            rounded_rect(sl, x, y, col_w, ROW_H, cell_fill,
                         corner=0.03, line=LINE_COLOR)
            label(sl, cell, x + 0.06, y + 0.05, col_w - 0.1, ROW_H - 0.08,
                  size=8, color=cell_col)

    # Bottom note
    rect(sl, 0.4, 7.65, 9.2, 0.03, fill=LINE_COLOR)
    label(sl, "现有架构 = 架构 0（已上线） · 架构 A/B/C = 待实现", 0.4, 7.72, 9.2,
          size=9, color=TEXT_DIM)
    return sl


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Implementation: Hook → mio-server
# ═══════════════════════════════════════════════════════════════════════════════
def slide_hook方案(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    rect(sl, 0, 0, 10, 0.06, fill=BRAND_LIME)

    label(sl, "06  Hook → mio-server 接入方案", 0.4, 0.25, 9, size=20,
          bold=True, color=TEXT_WHITE)
    subtitle(sl, "把本地 CC hook 的数据打入 mio-server，有 3 种实现路径", top=0.72)

    options = [
        (
            "方案 1：HTTP API 直推",
            ACCENT_BLUE,
            [
                "hook socket 改为 HTTP POST",
                "POST /v1/sessions/:id/messages",
                "需要 auth token（deviceId JWT）",
                "支持批推，消息去重（localId）",
                "实现难度：★☆☆☆☆（最简单）",
            ]
        ),
        (
            "方案 2：Socket.io 客户端",
            BRAND_LIME,
            [
                "hook 实现 Socket.io Client",
                "连接到 mio-server /v1/updates",
                "实时性好，支持 ack",
                "需要 WebSocket 支持（部分环境受限）",
                "实现难度：★★☆☆☆",
            ]
        ),
        (
            "方案 3：Webhook 中转站",
            ACCENT_RED,
            [
                "本地保留 hook socket",
                "MioIsland 转发到 mio-server",
                "兼容现有架构，改动最小",
                "notch 需要在线（Mac 要开）",
                "实现难度：★☆☆☆☆",
            ]
        ),
    ]

    for i, (title, color, bullets) in enumerate(options):
        x = 0.4 + i * 3.15
        # Card
        rounded_rect(sl, x, 1.3, 3.0, 5.5, RGBColor(0x0D, 0x0D, 0x0D),
                     corner=0.1, line=color)
        # Header
        rect(sl, x, 1.3, 3.0, 0.5, fill=color)
        label(sl, title, x, 1.32, 3.0, 0.46, size=11, bold=True,
              color=BRAND_DARK)
        # Bullets
        txb = sl.shapes.add_textbox(
            Inches(x + 0.15), Inches(1.9), Inches(2.75), Inches(4.7))
        tf = txb.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            # Color last bullet differently
            if "★" in b:
                run.font.color.rgb = color
                run.font.bold = True
            else:
                run.font.color.rgb = TEXT_DIM
            run.text = f"• {b}"
            run.font.size = Pt(9.5)

    # Bottom recommendation
    rect(sl, 0.4, 7.0, 9.2, 0.9, fill=RGBColor(0x0A, 0x20, 0x0A),
         line=BRAND_LIME)
    label(sl, "推荐路径：先用方案 1（HTTP 直推）快速验证 → 再迁移到方案 2（Socket.io）提升实时性",
          0.6, 7.1, 8.8, size=10.5, bold=True, color=BRAND_LIME)
    label(sl, "mio-server 的 /v1/sessions/:id/messages 已支持 batch POST + localId 去重",
          0.6, 7.45, 8.8, size=9, color=TEXT_DIM)
    return sl


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Session Lifecycle
# ═══════════════════════════════════════════════════════════════════════════════
def slide_lifecycle(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    rect(sl, 0, 0, 10, 0.06, fill=BRAND_LIME)

    label(sl, "07  Session 生命周期", 0.4, 0.25, 9, size=22, bold=True,
          color=TEXT_WHITE)
    subtitle(sl, "一条消息从 CC 产生到 iPhone 显示的完整路径", top=0.75)

    STEPS = [
        ("1", "CC 执行命令", "Claude Code 运行，产生 assistant / user / tool 消息"),
        ("2", "Hook 捕获", "PreToolUse / PostToolUse / Stop hooks 读取 JSONL event"),
        ("3", "发送到 mio-server", "HTTP POST /v1/sessions/:id/messages 或 Socket.io emit"),
        ("4", "Server 存储", "db.sessionMessage 表，seq 自增，localId 去重"),
        ("5", "Socket.io 广播", "eventRouter.emitUpdate → 所有 linked devices"),
        ("6", "iPhone 轮询", "GET /v1/sessions/:id/messages?after_seq=N 获取新消息"),
        ("7", "iPhone 渲染", "CodeLight App 显示 session 列表 + 消息内容"),
    ]

    for i, (num, title, desc) in enumerate(STEPS):
        y = 1.25 + i * 0.92
        # Number circle
        col = BRAND_LIME if i % 2 == 0 else ACCENT_BLUE
        circ = rounded_rect(sl, 0.5, y, 0.5, 0.5, fill=col, corner=0.25)
        label(sl, num, 0.5, y + 0.06, 0.5, 0.38, size=14, bold=True,
              color=BRAND_DARK)
        # Content card
        rounded_rect(sl, 1.1, y, 8.4, 0.78, RGBColor(0x0D, 0x0D, 0x0D),
                     corner=0.05, line=LINE_COLOR)
        label(sl, title, 1.2, y + 0.05, 3, 0.3, size=10, bold=True,
              color=TEXT_WHITE)
        label(sl, desc, 1.2, y + 0.36, 8.2, 0.36, size=9, color=TEXT_DIM)
        # Arrow
        if i < len(STEPS) - 1:
            solid_line(sl, 0.75, y + 0.55, 0.75, y + 0.9,
                       color=LINE_COLOR, w=0.8)

    return sl


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Next Steps
# ═══════════════════════════════════════════════════════════════════════════════
def slide_next(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 10, 8.5, fill=BG_DARK)
    rect(sl, 0, 0, 10, 0.06, fill=BRAND_LIME)

    label(sl, "08  下一步行动", 0.4, 0.25, 9, size=22, bold=True,
          color=TEXT_WHITE)
    subtitle(sl, "从零到一，渐进式实现 Server CC 模式", top=0.75)

    phases = [
        ("Phase 0", "验证现有架构", BRAND_LIME, [
            "Mac notch + mio-server 配对成功",
            "iPhone 可以看到本地 CC session",
            "APNs 通知正常",
        ]),
        ("Phase 1", "HTTP Hook 直推（最小MVP）", ACCENT_BLUE, [
            "改造 hook socket → HTTP POST",
            "POST /v1/sessions/:id/messages",
            "Server CC 机器持有 Mac deviceId JWT",
            "iPhone 轮询验证 session 出现",
        ]),
        ("Phase 2", "Socket.io 实时推送", BRAND_LIME, [
            "升级为 Socket.io 客户端",
            "支持 ack + 实时 phase 通知",
            "支持 APNs completion/approval 推送",
        ]),
        ("Phase 3", "独立 DeviceId 体系（可选）", TEXT_DIM, [
            "Server CC 拥有自己的 deviceId",
            "新建 DeviceLink 表（server CC ↔ mio-server）",
            "iPhone 配对新增 'Server CC' 类型",
        ]),
    ]

    for i, (phase, title, color, items) in enumerate(phases):
        x = 0.4 + i * 2.35
        rounded_rect(sl, x, 1.3, 2.2, 5.8,
                     RGBColor(0x0D, 0x0D, 0x0D), corner=0.1, line=color)
        # Phase tag
        rect(sl, x, 1.3, 2.2, 0.45, fill=color)
        label(sl, phase, x, 1.32, 2.2, 0.42, size=10, bold=True,
              color=BRAND_DARK)
        # Title
        label(sl, title, x + 0.1, 1.85, 2.0, 0.6, size=9.5, bold=True,
              color=TEXT_WHITE, wrap=True)
        # Items
        txb = sl.shapes.add_textbox(
            Inches(x + 0.1), Inches(2.55), Inches(2.05), Inches(4.3))
        tf = txb.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = f"• {item}"
            run.font.size = Pt(8.5)
            run.font.color.rgb = color if color != TEXT_DIM else TEXT_DIM
        # Arrow between phases
        if i < len(phases) - 1:
            label(sl, "→", x + 2.2, 3.8, 0.15, size=18, bold=True,
                  color=LINE_COLOR)

    rect(sl, 0.4, 7.25, 9.2, 0.7, fill=RGBColor(0x0A, 0x20, 0x0A),
         line=BRAND_LIME)
    label(sl, "关键前提：Server CC 机器需要持有有效的 Mac deviceId JWT（从已配对的 Mac 导出，或新建 deviceId）",
          0.6, 7.32, 8.8, size=10, color=BRAND_LIME)
    return sl


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(8.5)

    cover(prs)
    slide_overview(prs)
    slide_arch_current(prs)
    slide_arch_server(prs)
    slide_arch_ssh(prs)
    slide_comparison(prs)
    slide_hook方案(prs)
    slide_lifecycle(prs)
    slide_next(prs)

    out = "/Users/toby/Documents/Projects/MioServer/mio-architecture-comparison.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
